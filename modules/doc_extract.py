"""
文档文本抽取模块 - 支持PDF、Word、PPT、Excel、TXT等格式
"""

from pathlib import Path

# 文档处理库（可选导入）
try:
    import fitz  # PyMuPDF for PDF
except ImportError:
    fitz = None
    print("警告: PyMuPDF未安装，PDF文本抽取功能不可用")

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None
    print("警告: python-docx未安装，Word文档处理功能不可用")

try:
    from pptx import Presentation
except ImportError:
    Presentation = None
    print("警告: python-pptx未安装，PPT文档处理功能不可用")

try:
    from openpyxl import load_workbook
except ImportError:
    load_workbook = None
    print("警告: openpyxl未安装，Excel处理功能不可用")


def extract_text_from_pdf(file_path: str) -> str:
    """从PDF抽取文本"""
    if fitz is None:
        return ""
    try:
        doc = fitz.open(file_path)
        text_parts = []
        for page in doc:
            text_parts.append(page.get_text())
        doc.close()
        return "\n".join(text_parts)
    except Exception as e:
        print(f"PDF抽取错误: {e}")
        return ""


def extract_text_from_docx(file_path: str) -> str:
    """从Word文档抽取文本"""
    if DocxDocument is None:
        return ""
    try:
        doc = DocxDocument(file_path)
        text_parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        # 也抽取表格中的文本
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    text_parts.append(" | ".join(row_text))
        return "\n".join(text_parts)
    except Exception as e:
        print(f"Word抽取错误: {e}")
        return ""


def extract_text_from_pptx(file_path: str) -> str:
    """从PPT抽取文本"""
    if Presentation is None:
        return ""
    try:
        prs = Presentation(file_path)
        text_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = [f"[幻灯片{slide_num}]"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text)
            if len(slide_texts) > 1:
                text_parts.append("\n".join(slide_texts))
        return "\n\n".join(text_parts)
    except Exception as e:
        print(f"PPT抽取错误: {e}")
        return ""


def extract_text_from_xlsx(file_path: str) -> str:
    """从Excel抽取文本"""
    if load_workbook is None:
        return ""
    try:
        wb = load_workbook(file_path, data_only=True)
        text_parts = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_texts = [f"[工作表: {sheet_name}]"]
            for row in sheet.iter_rows(values_only=True):
                row_values = [str(cell) for cell in row if cell is not None]
                if row_values:
                    sheet_texts.append(" | ".join(row_values))
            if len(sheet_texts) > 1:
                text_parts.append("\n".join(sheet_texts))
        wb.close()
        return "\n\n".join(text_parts)
    except Exception as e:
        print(f"Excel抽取错误: {e}")
        return ""


def extract_text_from_document(file_path: str, filename: str) -> str:
    """根据文件类型抽取文本"""
    ext = Path(filename).suffix.lower()
    
    if ext == '.pdf':
        return extract_text_from_pdf(file_path)
    elif ext in ['.docx', '.doc']:
        return extract_text_from_docx(file_path)
    elif ext in ['.pptx', '.ppt']:
        return extract_text_from_pptx(file_path)
    elif ext in ['.xlsx', '.xls']:
        return extract_text_from_xlsx(file_path)
    elif ext == '.txt':
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        except Exception as e:
            print(f"TXT读取错误: {e}")
            return ""
    else:
        return ""


def extract_table_from_file(file_path: str, filename: str) -> str:
    """
    从Excel或CSV文件抽取表格内容，格式化为适合AI理解的文本
    限制最大3000字符，避免prompt过长
    """
    ext = Path(filename).suffix.lower()
    MAX_TABLE_LENGTH = 3000
    
    if ext == '.csv':
        try:
            import csv
            rows = []
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                reader = csv.reader(f)
                for row in reader:
                    if any(cell.strip() for cell in row):
                        rows.append(" | ".join(row))
                    if len("\n".join(rows)) > MAX_TABLE_LENGTH:
                        break
            result = "\n".join(rows)
            if len(result) > MAX_TABLE_LENGTH:
                result = result[:MAX_TABLE_LENGTH] + "\n...(数据过长，已截取)"
            return result
        except Exception as e:
            print(f"CSV读取错误: {e}")
            return ""
    
    elif ext in ['.xlsx', '.xls']:
        if load_workbook is None:
            return ""
        try:
            wb = load_workbook(file_path, data_only=True)
            all_tables = []
            total_length = 0
            
            for sheet_name in wb.sheetnames:
                if total_length > MAX_TABLE_LENGTH:
                    break
                    
                sheet = wb[sheet_name]
                table_rows = []
                
                for row in sheet.iter_rows(values_only=True):
                    row_values = []
                    for cell in row:
                        if cell is not None:
                            row_values.append(str(cell))
                        else:
                            row_values.append("")
                    # 只添加非空行
                    if any(v.strip() for v in row_values):
                        table_rows.append(" | ".join(row_values))
                    
                    # 检查长度
                    if len("\n".join(table_rows)) > MAX_TABLE_LENGTH:
                        break
                
                if table_rows:
                    if len(wb.sheetnames) > 1:
                        all_tables.append(f"[{sheet_name}]\n" + "\n".join(table_rows))
                    else:
                        all_tables.append("\n".join(table_rows))
                    total_length = len("\n\n".join(all_tables))
            
            wb.close()
            result = "\n\n".join(all_tables)
            if len(result) > MAX_TABLE_LENGTH:
                result = result[:MAX_TABLE_LENGTH] + "\n...(数据过长，已截取)"
            return result
        except Exception as e:
            print(f"Excel表格读取错误: {e}")
            return ""
    
    return ""
